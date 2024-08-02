from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os
import io

# ファイルパスを絶対パスで指定
input_pptx = '/Users/yn_mbp/git/yn-scripts/master.pptx'
output_pptx = '/Users/yn_mbp/git/yn-scripts/output_presentation.pptx'

# プレゼンテーションオブジェクトを作成
prs = Presentation(input_pptx)

# 一時的なJPEG保存フォルダ
temp_folder = 'temp_jpegs'
os.makedirs(temp_folder, exist_ok=True)

def convert_to_jpeg(image_data, output_path):
    img = Image.open(io.BytesIO(image_data))
    rgb_img = img.convert('RGB')
    rgb_img.save(output_path, 'JPEG')

# 各スライドをループ
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            image = shape.image
            if image.ext == 'png':
                png_data = image.blob  # This gets the image data as bytes
                jpeg_path = os.path.join(temp_folder, f'{os.path.basename(shape.name)}.jpeg')
                convert_to_jpeg(png_data, jpeg_path)
                
                # 元の画像を削除
                shape._element.getparent().remove(shape._element)
                
                # 新しいJPEG画像を挿入
                slide.shapes.add_picture(jpeg_path, shape.left, shape.top, shape.width, shape.height)

# PowerPointファイルとして保存
prs.save(output_pptx)

# 一時フォルダの削除
import shutil
shutil.rmtree(temp_folder)